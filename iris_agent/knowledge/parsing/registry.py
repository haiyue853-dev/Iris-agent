"""Suffix-keyed parser registry: one function per source format."""

from __future__ import annotations

from io import BytesIO
from typing import Callable

from iris_agent.knowledge.parsing.base import ParsingError, ParsedSection

_MAX_SECTION_CHARS = 60_000


def supported_suffixes() -> tuple[str, ...]:
    return tuple(sorted(_PARSERS))


def _clip(text: str) -> str:
    return text[:_MAX_SECTION_CHARS]


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8", "gb18030"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _parse_text(content: bytes, *, name: str) -> tuple[ParsedSection, ...]:
    text = _decode_text(content).strip()
    if not text:
        raise ParsingError("文档未包含可提取文本")
    return (ParsedSection(_clip(text)),)


ImageDescriber = Callable[[bytes, str], str]


def _parse_pdf(content: bytes, *, name: str, image_describer: ImageDescriber | None = None) -> tuple[ParsedSection, ...]:
    import pymupdf as fitz

    sections: list[ParsedSection] = []
    with fitz.open(stream=content, filetype="pdf") as document:
        for page_number, page in enumerate(document, start=1):
            location = f"第 {page_number} 页"
            text = page.get_text("text").strip()
            if text:
                sections.append(ParsedSection(_clip(text), location=location))
            try:
                tables = page.find_tables()
            except Exception:
                tables = None
            for table in (tables.tables if tables is not None else []):
                try:
                    rows = table.extract()
                except Exception:
                    continue
                lines = ["\t".join(str(cell or "").strip() for cell in row) for row in rows]
                table_text = "\n".join(line for line in lines if line.strip()).strip()
                if table_text:
                    sections.append(ParsedSection(_clip(table_text), kind="table", location=location))
            if not text and image_describer is not None:
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                described = image_describer(pixmap.tobytes("png"), f"{name} · {location}").strip()
                if described:
                    sections.append(ParsedSection(_clip(described), kind="image", location=location))
    if not sections:
        raise ParsingError("PDF 未包含可提取文本（扫描件请先 OCR）")
    return tuple(sections)


def _parse_docx(content: bytes, *, name: str) -> tuple[ParsedSection, ...]:
    from docx import Document

    document = Document(BytesIO(content))
    sections: list[ParsedSection] = []
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    if paragraphs:
        sections.append(ParsedSection(_clip("\n".join(paragraphs))))
    for index, table in enumerate(document.tables, start=1):
        lines = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        table_text = "\n".join(line for line in lines if line.strip()).strip()
        if table_text:
            sections.append(ParsedSection(_clip(table_text), kind="table", location=f"表格 {index}"))
    if not sections:
        raise ParsingError("Word 文档未包含可提取文本")
    return tuple(sections)


def _parse_xlsx(content: bytes, *, name: str) -> tuple[ParsedSection, ...]:
    from openpyxl import load_workbook

    workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
    sections: list[ParsedSection] = []
    for sheet in workbook.worksheets:
        lines = [
            "\t".join("" if cell is None else str(cell).strip() for cell in row)
            for row in sheet.iter_rows(values_only=True)
        ]
        sheet_text = "\n".join(line for line in lines if line.strip()).strip()
        if sheet_text:
            sections.append(ParsedSection(_clip(sheet_text), kind="table", location=f"工作表 {sheet.title}"))
    workbook.close()
    if not sections:
        raise ParsingError("Excel 文档未包含可提取文本")
    return tuple(sections)


def _parse_xls(content: bytes, *, name: str) -> tuple[ParsedSection, ...]:
    import xlrd

    workbook = xlrd.open_workbook(file_contents=content)
    sections: list[ParsedSection] = []
    for sheet in workbook.sheets():
        lines = []
        for row_index in range(sheet.nrows):
            line = "\t".join(str(sheet.cell_value(row_index, col)).strip() for col in range(sheet.ncols))
            if line.strip():
                lines.append(line)
        if lines:
            sections.append(ParsedSection(_clip("\n".join(lines)), kind="table", location=f"工作表 {sheet.name}"))
    if not sections:
        raise ParsingError("Excel 文档未包含可提取文本")
    return tuple(sections)


def _parse_pptx(content: bytes, *, name: str) -> tuple[ParsedSection, ...]:
    from pptx import Presentation

    presentation = Presentation(BytesIO(content))
    sections: list[ParsedSection] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if text:
                    lines.append(text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(notes)
        if lines:
            sections.append(ParsedSection(_clip("\n".join(lines)), location=f"幻灯片 {index}"))
    if not sections:
        raise ParsingError("演示文稿未包含可提取文本")
    return tuple(sections)


def _parse_html(content: bytes, *, name: str) -> tuple[ParsedSection, ...]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode_text(content), "html.parser")
    for tag in soup(["script", "style", "noscript", "head"]):
        tag.decompose()
    text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip()).strip()
    if not text:
        raise ParsingError("网页未包含可提取文本")
    return (ParsedSection(_clip(text)),)


def _parse_mhtml(content: bytes, *, name: str) -> tuple[ParsedSection, ...]:
    import email
    from email import policy

    message = email.message_from_bytes(content, policy=policy.default)
    payload: bytes | None = None
    for part in message.walk():
        if part.get_content_type() == "text/html":
            candidate = part.get_payload(decode=True)
            if candidate:
                payload = candidate
                break
    if payload is None:
        return _parse_html(content, name=name)
    return _parse_html(payload, name=name)


def _parse_image(content: bytes, *, name: str, image_describer: ImageDescriber | None = None) -> tuple[ParsedSection, ...]:
    if image_describer is None:
        raise ParsingError("图片知识入库需要启用视觉解析模型")
    text = image_describer(content, name).strip()
    if not text:
        raise ParsingError("视觉解析模型未返回可索引内容")
    return (ParsedSection(_clip(text), kind="image"),)


_PARSERS = {
    ".txt": _parse_text,
    ".md": _parse_text,
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".xlsx": _parse_xlsx,
    ".xls": _parse_xls,
    ".pptx": _parse_pptx,
    ".html": _parse_html,
    ".htm": _parse_html,
    ".mhtml": _parse_mhtml,
    ".png": _parse_image,
    ".jpg": _parse_image,
    ".jpeg": _parse_image,
    ".webp": _parse_image,
}


def parse_document(suffix: str, content: bytes, *, name: str = "", image_describer: ImageDescriber | None = None):
    """Parse raw bytes into ordered text sections; raise ParsingError when impossible."""
    from iris_agent.knowledge.parsing.base import ParsedDocument

    if not isinstance(suffix, str) or not isinstance(content, (bytes, bytearray)):
        raise ValueError("suffix and content are required")
    parser = _PARSERS.get(suffix.lower())
    if parser is None:
        raise ParsingError(f"不支持的知识库文件类型：{suffix or '未知'}")
    try:
        if parser in {_parse_pdf, _parse_image}:
            sections = parser(bytes(content), name=name, image_describer=image_describer)
        else:
            sections = parser(bytes(content), name=name)
    except ParsingError:
        raise
    except ImportError as exc:
        raise ParsingError(f"缺少解析依赖：{exc}") from exc
    except Exception as exc:
        raise ParsingError(f"无法解析{name or '文档'}：{exc}") from exc
    warnings: list[str] = []
    if any(section.text.endswith("…") or len(section.text) >= _MAX_SECTION_CHARS for section in sections):
        warnings.append("部分段落超过单段长度上限，已截断")
    return ParsedDocument(sections=tuple(sections), warnings=tuple(warnings))
